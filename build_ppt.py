"""
NEXUS DV2.0 — Presentation Builder
Generates a professional PowerPoint deck for the NEXUS platform.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── Color palette ──────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x14, 0x2B, 0x4C)   # deep navy — primary background
MID_BLUE   = RGBColor(0x1A, 0x5C, 0x9A)   # medium blue — section headers
TEAL       = RGBColor(0x0D, 0x8A, 0x84)   # teal — accent
SLATE      = RGBColor(0x3A, 0x4A, 0x60)   # slate — secondary boxes
ORANGE     = RGBColor(0xE3, 0x6C, 0x1A)   # orange — highlight/callout
GREEN      = RGBColor(0x19, 0x8A, 0x55)   # green — approval/positive
RED        = RGBColor(0xC0, 0x39, 0x2B)   # red — reject/negative
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)   # light blue — box fill
LIGHT_TEAL = RGBColor(0xD1, 0xF2, 0xEB)   # light teal — box fill
LIGHT_GREY = RGBColor(0xF2, 0xF4, 0xF7)   # light grey — box fill
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF8, 0xF9, 0xFB)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
MID_TEXT   = RGBColor(0x3D, 0x4D, 0x6A)
AMBER      = RGBColor(0xF3, 0x9C, 0x12)   # amber — warning/notes

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helper utilities ──────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=Pt(0.75), radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    fill_fmt = shape.fill
    if fill:
        fill_fmt.solid()
        fill_fmt.fore_color.rgb = fill
    else:
        fill_fmt.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = line_width
    else:
        line.fill.background()
    return shape


def add_rounded_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=Pt(0.75)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        5,  # ROUNDED_RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    # Set corner radius
    shape.adjustments[0] = 0.05
    fill_fmt = shape.fill
    if fill:
        fill_fmt.solid()
        fill_fmt.fore_color.rgb = fill
    else:
        fill_fmt.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = line_width
    else:
        line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=Pt(12), bold=False, color=DARK_TEXT,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def txt_box(slide, text, x, y, w, h,
            fill=None, line_color=None,
            size=Pt(11), bold=False, text_color=DARK_TEXT,
            align=PP_ALIGN.CENTER, v_anchor=None, italic=False, line_width=Pt(0.75)):
    """Rounded rect with centred text."""
    shape = add_rounded_rect(slide, x, y, w, h, fill=fill,
                             line_color=line_color, line_width=line_width)
    tf = shape.text_frame
    tf.word_wrap = True
    from pptx.enum.text import MSO_ANCHOR
    if v_anchor:
        tf.auto_size = None
        tf.vertical_anchor = v_anchor
    else:
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = text_color
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=MID_BLUE, width=Pt(1.5)):
    """Simple line connector drawn as a rectangle-based arrow using a line shape."""
    from pptx.util import Inches
    connector = slide.shapes.add_connector(
        1,  # STRAIGHT
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def slide_header(slide, title, subtitle=None,
                 header_h=1.05, bg_color=NAVY):
    # Full-width header bar
    add_rect(slide, 0, 0, 13.33, header_h, fill=bg_color)
    # Title
    txt(slide, title, 0.35, 0.12, 11.5, 0.6,
        size=Pt(24), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.68, 11.5, 0.35,
            size=Pt(12), bold=False, color=RGBColor(0xA8, 0xC7, 0xE8),
            align=PP_ALIGN.LEFT)
    # Thin accent line below header
    add_rect(slide, 0, header_h, 13.33, 0.04, fill=TEAL)


def footer(slide, page_num):
    add_rect(slide, 0, 7.25, 13.33, 0.25, fill=RGBColor(0x0D, 0x1B, 0x35))
    txt(slide, "NEXUS DV2.0  |  AI-Assisted Data Vault Automation  |  Confidential",
        0.3, 7.27, 10, 0.2, size=Pt(8), color=RGBColor(0x80, 0x9A, 0xBF))
    txt(slide, str(page_num), 12.8, 7.27, 0.4, 0.2,
        size=Pt(8), color=RGBColor(0x80, 0x9A, 0xBF), align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank_layout)

# Full dark background
add_rect(slide, 0, 0, 13.33, 7.5, fill=NAVY)

# Decorative accent bar (left vertical)
add_rect(slide, 0, 0, 0.18, 7.5, fill=TEAL)

# Decorative horizontal accent
add_rect(slide, 0.18, 3.55, 13.15, 0.06, fill=TEAL)

# Diagonal geometric accent (top right triangle-ish using rectangles)
add_rect(slide, 9.5, 0, 3.83, 3.4, fill=RGBColor(0x1A, 0x3A, 0x6B))
add_rect(slide, 10.5, 0, 2.83, 2.2, fill=RGBColor(0x1E, 0x44, 0x7A))

# Main title
txt(slide, "NEXUS", 0.6, 1.0, 9, 1.5,
    size=Pt(72), bold=True, color=WHITE)
txt(slide, "DV 2.0", 0.6, 2.3, 9, 1.2,
    size=Pt(52), bold=True, color=TEAL)

# Subtitle
txt(slide, "AI-Assisted Data Vault Automation Platform", 0.6, 3.75, 10, 0.6,
    size=Pt(20), color=RGBColor(0xA8, 0xC7, 0xE8))

# Tag line
txt(slide, "Built entirely on Snowflake Cortex  ·  Human-in-the-Loop Design  ·  Enterprise Banking",
    0.6, 4.4, 11, 0.4,
    size=Pt(13), italic=True, color=RGBColor(0x70, 0x8A, 0xB0))

# Bottom bar
add_rect(slide, 0, 6.95, 13.33, 0.55, fill=RGBColor(0x0A, 0x16, 0x2B))
txt(slide, "Confidential  |  Enterprise Data Architecture", 0.6, 7.0, 10, 0.4,
    size=Pt(10), color=RGBColor(0x60, 0x7A, 0xA0))

# Snowflake logo placeholder — snowflake symbol text
txt(slide, "❄", 11.8, 6.9, 1.3, 0.6,
    size=Pt(28), color=RGBColor(0x29, 0xB5, 0xE8), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Intent & Design Philosophy
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
slide_header(slide, "Intent & Design Philosophy",
             "What NEXUS is built to solve — and how we keep humans in control")
footer(slide, 2)

# Three pillars
pillars = [
    ("🤖  AI Accelerates", MID_BLUE,
     "Claude Opus 4.6 analyzes source table structure, column statistics, and existing vault "
     "patterns to propose hubs, links, and satellites — in seconds rather than days."),
    ("👤  Modelers Decide", TEAL,
     "Every AI proposal lands in a Design Workbench for review. Modelers can edit names, "
     "reassign columns, add notes, split satellites, or reject entirely before anything is written "
     "to the registry."),
    ("✅  Governance Enforces", GREEN,
     "No entity enters the vault registry without an explicit APPROVE action. Every decision is "
     "captured in an immutable audit log — who generated, who edited, who approved, and when."),
]

for i, (title, color, body) in enumerate(pillars):
    bx = 0.3 + i * 4.35
    add_rounded_rect(slide, bx, 1.25, 4.1, 4.9,
                     fill=WHITE, line_color=color, line_width=Pt(2))
    # Color top cap
    shape = slide.shapes.add_shape(5, Inches(bx), Inches(1.25), Inches(4.1), Inches(0.7))
    shape.adjustments[0] = 0.04
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    txt(slide, title, bx + 0.15, 1.32, 3.8, 0.55,
        size=Pt(15), bold=True, color=WHITE)
    txt(slide, body, bx + 0.2, 2.05, 3.7, 3.8,
        size=Pt(11.5), color=MID_TEXT, wrap=True)

# Bottom callout
add_rounded_rect(slide, 0.3, 6.3, 12.73, 0.65,
                 fill=RGBColor(0xFF, 0xF3, 0xE0), line_color=ORANGE, line_width=Pt(1.5))
txt(slide, "💡  Core principle:  AI provides its best guess with explicit confidence ratings. "
    "Modeler overrides always win. The system prompt, abbreviation table, profiling stats, "
    "and modeler notes are all passed to the AI as structured context — nothing is left to hallucination.",
    0.5, 6.35, 12.4, 0.55, size=Pt(10.5), color=RGBColor(0x7D, 0x3C, 0x00))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Human-in-the-Loop Workflow
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
slide_header(slide, "Human-in-the-Loop Workflow",
             "Every AI proposal passes through explicit modeler review and approval gates")
footer(slide, 3)

# ── Flow steps ────────────────────────────────────────────────────────────────
steps = [
    ("1\nIdentify\nSource", MID_BLUE, "Browse any\nSnowflake DB"),
    ("2\nProfile\nTable",   TEAL,     "Column stats,\nPK scoring"),
    ("3\nReview\nStats",    SLATE,    "Confirm BK,\nchange freq"),
    ("4\nAI Generate\nProposal", ORANGE, "Cortex call\nw/ full context"),
    ("5\nDesign\nWorkbench", MID_BLUE, "Edit entities\n& columns"),
    ("6\nApprove /\nReject", GREEN,   "Gate before\nregistry write"),
    ("7\nWrite to\nRegistry", TEAL,   "DDL +\ndbt ready"),
]

box_w = 1.55
box_h = 1.1
start_x = 0.28
y_step  = 1.35
gap     = 0.12

for i, (label, color, sub) in enumerate(steps):
    bx = start_x + i * (box_w + gap)
    # Main box
    shape = slide.shapes.add_shape(5, Inches(bx), Inches(1.45), Inches(box_w), Inches(box_h))
    shape.adjustments[0] = 0.06
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    txt(slide, label, bx + 0.05, 1.5, box_w - 0.1, box_h - 0.05,
        size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Sub-label below
    txt(slide, sub, bx + 0.05, 2.62, box_w - 0.1, 0.5,
        size=Pt(9), color=MID_TEXT, align=PP_ALIGN.CENTER)
    # Arrow (not after last)
    if i < len(steps) - 1:
        ax = bx + box_w + 0.01
        add_arrow(slide, ax, 1.45 + box_h / 2, ax + gap + 0.01, 1.45 + box_h / 2,
                  color=MID_BLUE, width=Pt(2))

# ── Decision diamond area ─────────────────────────────────────────────────────
# Show what happens at step 6 — approve path and reject path
add_rect(slide, 0.3, 3.25, 12.73, 0.04, fill=RGBColor(0xCC, 0xD8, 0xE8))

# Approve branch
add_rounded_rect(slide, 0.3, 3.38, 5.9, 0.75,
                 fill=RGBColor(0xE8, 0xF8, 0xF0), line_color=GREEN, line_width=Pt(1.5))
txt(slide, "✅  APPROVE  →  Entity written to META.DV_ENTITY + DV_ENTITY_COLUMN  →  Available for dbt generation",
    0.5, 3.42, 5.5, 0.65, size=Pt(10.5), color=RGBColor(0x0A, 0x5C, 0x35))

# Reject branch
add_rounded_rect(slide, 6.5, 3.38, 3.3, 0.75,
                 fill=RGBColor(0xFD, 0xED, 0xEC), line_color=RED, line_width=Pt(1.5))
txt(slide, "❌  REJECT  →  Proposal archived,\nno registry change",
    6.65, 3.42, 3.0, 0.65, size=Pt(10.5), color=RGBColor(0x7D, 0x0A, 0x0A))

# Re-generate branch
add_rounded_rect(slide, 10.1, 3.38, 2.93, 0.75,
                 fill=RGBColor(0xFF, 0xF3, 0xE0), line_color=AMBER, line_width=Pt(1.5))
txt(slide, "🔄  RE-RUN  →  Add modeler\nnotes, regenerate",
    10.25, 3.42, 2.65, 0.65, size=Pt(10.5), color=RGBColor(0x7D, 0x4A, 0x00))

# ── Key controls table ────────────────────────────────────────────────────────
add_rounded_rect(slide, 0.3, 4.3, 12.73, 2.65,
                 fill=WHITE, line_color=RGBColor(0xCC, 0xD6, 0xE8), line_width=Pt(1))

txt(slide, "Key Human Control Points", 0.5, 4.35, 6, 0.35,
    size=Pt(13), bold=True, color=NAVY)

controls = [
    ("Modeler Notes Override",     "Free-text field that overrides ALL statistical PK detection. AI explicitly instructed to cite these notes in every entity rationale."),
    ("PK Confirmation",            "Modeler selects or confirms the business key from ranked PK candidates before AI proposal is generated — injected as a 100-score candidate."),
    ("Column Definitions",         "Approved definitions stored in META.DV_COLUMN_DEFINITIONS are passed verbatim to AI — no paraphrasing allowed."),
    ("Workspace Versioning",       "Every edit creates a new version. PARENT_WORKSPACE tracks lineage. Status: DRAFT → IN_REVIEW → APPROVED / SUPERSEDED."),
    ("Immutable Audit Log",        "Every action (PROFILE, GENERATE, SAVE, APPROVE, EXPORT) written to META.DV_AUDIT_LOG with timestamp and user."),
]

for i, (ctrl, desc) in enumerate(controls):
    cy = 4.78 + i * 0.42
    txt(slide, f"▸  {ctrl}", 0.5, cy, 3.1, 0.38, size=Pt(10), bold=True, color=NAVY)
    txt(slide, desc, 3.65, cy, 9.1, 0.38, size=Pt(10), color=MID_TEXT)
    if i < len(controls) - 1:
        add_rect(slide, 0.4, cy + 0.38, 12.5, 0.02, fill=LIGHT_GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — High-Level Architecture
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
slide_header(slide, "High-Level Architecture",
             "All components run inside Snowflake — no external tools required")
footer(slide, 4)

# ── Outer Snowflake boundary ──────────────────────────────────────────────────
add_rounded_rect(slide, 0.22, 1.2, 12.89, 5.85,
                 fill=RGBColor(0xF0, 0xF5, 0xFF), line_color=RGBColor(0x29, 0xB5, 0xE8),
                 line_width=Pt(2))
txt(slide, "❄  Snowflake Platform Boundary", 0.4, 1.24, 5, 0.32,
    size=Pt(10), bold=True, color=RGBColor(0x29, 0xB5, 0xE8))

# ── Layer 0: Source Systems (outside boundary, left) ─────────────────────────
add_rounded_rect(slide, 0.28, 1.62, 1.72, 5.22,
                 fill=RGBColor(0xEB, 0xEE, 0xF4), line_color=SLATE, line_width=Pt(1.5))
txt(slide, "SOURCE\nSYSTEMS", 0.35, 1.65, 1.6, 0.55,
    size=Pt(9.5), bold=True, color=SLATE, align=PP_ALIGN.CENTER)

src_systems = ["ACCT_SYS", "CRM_SYS", "LOAN_SYS", "REF Data"]
for i, s in enumerate(src_systems):
    sy = 2.35 + i * 1.05
    add_rounded_rect(slide, 0.38, sy, 1.52, 0.62,
                     fill=SLATE, line_color=None)
    txt(slide, s, 0.38, sy + 0.08, 1.52, 0.45,
        size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow from sources into NEXUS
add_arrow(slide, 1.9, 3.75, 2.4, 3.75, color=SLATE, width=Pt(2))
txt(slide, "any accessible\nSnowflake role", 1.83, 3.85, 1.4, 0.4,
    size=Pt(7.5), italic=True, color=SLATE, align=PP_ALIGN.CENTER)

# ── Layer 1: META Schema (main NEXUS area) ─────────────────────────────────────
# Phase 1 — Registry & Foundation
add_rounded_rect(slide, 2.45, 1.58, 3.55, 2.2,
                 fill=WHITE, line_color=MID_BLUE, line_width=Pt(1.5))
shape = slide.shapes.add_shape(5, Inches(2.45), Inches(1.58), Inches(3.55), Inches(0.52))
shape.adjustments[0] = 0.04
shape.fill.solid(); shape.fill.fore_color.rgb = MID_BLUE
shape.line.fill.background()
txt(slide, "Phase 1 · Registry & Foundation", 2.55, 1.62, 3.3, 0.42,
    size=Pt(10), bold=True, color=WHITE)

registry_items = [
    "META.DV_ENTITY  &  DV_ENTITY_COLUMN",
    "META.DV_HASH_DEFINITION",
    "META.DV_SOURCE_MAPPING",
    "META.DV_ABBREVIATION  (Erwin .ABR sync)",
    "META.DV_AI_SYSTEM_PROMPT  (versioned)",
    "META.DV_AUDIT_LOG",
]
for i, item in enumerate(registry_items):
    txt(slide, f"· {item}", 2.55, 2.18 + i * 0.25, 3.3, 0.25,
        size=Pt(8.5), color=MID_TEXT)

# Phase 2 — Profiling Engine
add_rounded_rect(slide, 2.45, 3.93, 3.55, 1.82,
                 fill=WHITE, line_color=TEAL, line_width=Pt(1.5))
shape = slide.shapes.add_shape(5, Inches(2.45), Inches(3.93), Inches(3.55), Inches(0.45))
shape.adjustments[0] = 0.04
shape.fill.solid(); shape.fill.fore_color.rgb = TEAL
shape.line.fill.background()
txt(slide, "Phase 2 · Profiling Engine", 2.55, 3.95, 3.3, 0.38,
    size=Pt(10), bold=True, color=WHITE)
for i, item in enumerate(["SP_PROFILE_TABLE  (EXACT / HLL)",
                           "SP_DETECT_PK_CANDIDATES",
                           "SP_DETECT_CHANGE_FREQUENCY",
                           "META.DV_PROFILING_RUN + RESULTS"]):
    txt(slide, f"· {item}", 2.55, 4.45 + i * 0.28, 3.3, 0.26,
        size=Pt(8.5), color=MID_TEXT)

# Arrows between phases (vertical)
add_arrow(slide, 4.22, 3.78, 4.22, 3.93, color=TEAL, width=Pt(1.5))

# Phase 3 — AI Design Engine
add_rounded_rect(slide, 6.2, 1.58, 3.6, 2.2,
                 fill=WHITE, line_color=ORANGE, line_width=Pt(1.5))
shape = slide.shapes.add_shape(5, Inches(6.2), Inches(1.58), Inches(3.6), Inches(0.52))
shape.adjustments[0] = 0.04
shape.fill.solid(); shape.fill.fore_color.rgb = ORANGE
shape.line.fill.background()
txt(slide, "Phase 3 · AI Design Engine", 6.3, 1.62, 3.4, 0.42,
    size=Pt(10), bold=True, color=WHITE)
for i, item in enumerate(["SP_GENERATE_DV_PROPOSAL",
                           "Cortex AI_COMPLETE (Claude Opus 4.6)",
                           "META.DV_DESIGN_PROPOSAL",
                           "META.DV_DESIGN_WORKSPACE  (versioned)",
                           "Modeler Notes  ·  Confidence flags"]):
    txt(slide, f"· {item}", 6.3, 2.18 + i * 0.28, 3.3, 0.26,
        size=Pt(8.5), color=MID_TEXT)

# Arrow from Phase 2 to Phase 3
add_arrow(slide, 6.0, 3.0, 6.2, 3.0, color=ORANGE, width=Pt(2))
txt(slide, "profiling\nresults", 5.88, 3.05, 0.88, 0.38,
    size=Pt(7.5), italic=True, color=ORANGE, align=PP_ALIGN.CENTER)

# Streamlit UI box
add_rounded_rect(slide, 6.2, 3.93, 3.6, 1.82,
                 fill=WHITE, line_color=MID_BLUE, line_width=Pt(1.5))
shape = slide.shapes.add_shape(5, Inches(6.2), Inches(3.93), Inches(3.6), Inches(0.45))
shape.adjustments[0] = 0.04
shape.fill.solid(); shape.fill.fore_color.rgb = MID_BLUE
shape.line.fill.background()
txt(slide, "Streamlit in Snowflake (UI)", 6.3, 3.95, 3.4, 0.38,
    size=Pt(10), bold=True, color=WHITE)
for i, item in enumerate(["Identify Source  ·  Profile & Review",
                           "Design Workbench  ·  Edit + Approve",
                           "ER Diagram (Mermaid/D3)",
                           "Generate dbt (Phase 4)"]):
    txt(slide, f"· {item}", 6.3, 4.45 + i * 0.28, 3.4, 0.26,
        size=Pt(8.5), color=MID_TEXT)

# Arrow Phase 3 → Streamlit
add_arrow(slide, 8.0, 3.78, 8.0, 3.93, color=MID_BLUE, width=Pt(1.5))

# Phase 4 — Output
add_rounded_rect(slide, 10.05, 1.58, 2.95, 5.17,
                 fill=WHITE, line_color=GREEN, line_width=Pt(1.5))
shape = slide.shapes.add_shape(5, Inches(10.05), Inches(1.58), Inches(2.95), Inches(0.52))
shape.adjustments[0] = 0.04
shape.fill.solid(); shape.fill.fore_color.rgb = GREEN
shape.line.fill.background()
txt(slide, "Phase 4 · Output", 10.15, 1.62, 2.75, 0.42,
    size=Pt(10), bold=True, color=WHITE)

outputs = [
    ("RAW_VAULT Schema", "HUBs, LNKs, SATs\nDDL generated from registry"),
    ("dbt Models", "Staging → Raw Vault\ntransformation models"),
    ("Hash Specs", "Reproducible SHA2-256\nhash key definitions"),
    ("Lineage", "Source→Vault column\nmapping (DV_SOURCE_MAPPING)"),
    ("ER Diagrams", "Interactive vault\nentity diagrams"),
]
for i, (title, desc) in enumerate(outputs):
    oy = 2.25 + i * 0.9
    add_rounded_rect(slide, 10.15, oy, 2.75, 0.78,
                     fill=RGBColor(0xEB, 0xF9, 0xEF), line_color=GREEN, line_width=Pt(0.5))
    txt(slide, title, 10.25, oy + 0.04, 2.55, 0.3,
        size=Pt(9), bold=True, color=GREEN)
    txt(slide, desc, 10.25, oy + 0.34, 2.55, 0.38,
        size=Pt(8), color=MID_TEXT)

# Arrow Streamlit → Phase 4
add_arrow(slide, 9.8, 4.1, 10.05, 4.1, color=GREEN, width=Pt(2))
txt(slide, "APPROVE", 9.62, 4.13, 1.0, 0.22,
    size=Pt(7.5), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Arrow Phase 1 → Phase 3 (registry context)
add_arrow(slide, 6.0, 2.5, 6.2, 2.5, color=MID_BLUE, width=Pt(1.5))
txt(slide, "registry\ncontext", 5.85, 2.55, 0.95, 0.38,
    size=Pt(7.5), italic=True, color=MID_BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Context Passed to AI_COMPLETE
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
slide_header(slide, "How Context is Assembled & Passed to AI_COMPLETE",
             "SP_GENERATE_DV_PROPOSAL builds a rich, structured prompt at runtime — every call is fully grounded")
footer(slide, 5)

# Central AI_COMPLETE box
add_rounded_rect(slide, 5.2, 2.7, 3.0, 1.5,
                 fill=ORANGE, line_color=None)
txt(slide, "SNOWFLAKE\nCORTEX\nAI_COMPLETE\n(Claude Opus 4.6)", 5.2, 2.75, 3.0, 1.35,
    size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Context inputs — left column
ctx_left = [
    ("1  System Prompt", MID_BLUE, LIGHT_BLUE,
     "Assembled at runtime from META.DV_AI_SYSTEM_PROMPT\n"
     "7 sections: ROLE · NAMING · HASH STANDARDS · METADATA\n"
     "COLUMNS · SATELLITE RULES · LINK RULES · RESPONSE FORMAT\n"
     "Versioned, editable by architecture team without code change"),
    ("2  Registry Context", TEAL, LIGHT_TEAL,
     "All APPROVED entities from META.DV_ENTITY + DV_ENTITY_COLUMN\n"
     "Enables hub reuse detection before creating new entities\n"
     "Format: ENTITY_ID [TYPE] — logical_name\n"
     "  Columns: col:role:datatype, ..."),
    ("3  Abbreviation Table", SLATE, LIGHT_GREY,
     "All active rows from META.DV_ABBREVIATION\n"
     "Maps physical abbreviations → logical names + domain\n"
     "Keeps column naming consistent with Erwin .ABR file\n"
     "Example:  CUST = Customer  (domain: PARTY)"),
]

for i, (title, border_col, fill_col, body) in enumerate(ctx_left):
    by = 1.3 + i * 1.85
    add_rounded_rect(slide, 0.25, by, 4.7, 1.65,
                     fill=fill_col, line_color=border_col, line_width=Pt(1.5))
    txt(slide, title, 0.4, by + 0.06, 4.4, 0.35, size=Pt(11), bold=True, color=border_col)
    txt(slide, body, 0.4, by + 0.42, 4.4, 1.15, size=Pt(8.5), color=MID_TEXT)
    # Arrow to center box
    add_arrow(slide, 4.95, by + 0.82, 5.2, 3.45, color=border_col, width=Pt(1.2))

# Context inputs — right column
ctx_right = [
    ("4  Profiling Results", ORANGE, RGBColor(0xFF, 0xF3, 0xE0),
     "Per-column stats from latest DV_PROFILING_RESULTS run:\n"
     "uniqueness %, null %, min/max/avg length, top 5 values\n"
     "inferred data type, pattern detection (UUID/EMAIL/CODE)\n"
     "change_frequency (FAST/SLOW/STATIC), PK candidate flag"),
    ("5  PK Candidates", GREEN, RGBColor(0xE8, 0xF8, 0xF0),
     "Ranked candidates from META.DV_PK_CANDIDATES\n"
     "Scoring 0-100: uniqueness, null %, naming heuristics\n"
     "MODELER_SELECTED flag = overrides all statistics\n"
     "Injected as: ← MODELER CONFIRMED in the prompt"),
    ("6  Modeler Notes\n    & Col. Definitions", AMBER, RGBColor(0xFF, 0xFB, 0xE6),
     "Free-text modeler instructions — placed FIRST in prompt\n"
     "Explicit override: PK, deprecated cols, relationship hints\n"
     "Approved definitions from DV_COLUMN_DEFINITIONS copied\n"
     "verbatim — AI forbidden from paraphrasing these"),
]

for i, (title, border_col, fill_col, body) in enumerate(ctx_right):
    by = 1.3 + i * 1.85
    add_rounded_rect(slide, 8.38, by, 4.7, 1.65,
                     fill=fill_col, line_color=border_col, line_width=Pt(1.5))
    txt(slide, title, 8.53, by + 0.06, 4.4, 0.45, size=Pt(11), bold=True, color=border_col)
    txt(slide, body, 8.53, by + 0.45, 4.4, 1.12, size=Pt(8.5), color=MID_TEXT)
    # Arrow to center box
    add_arrow(slide, 8.38, by + 0.82, 8.2, 3.45, color=border_col, width=Pt(1.2))

# Output arrow
add_arrow(slide, 6.7, 4.2, 6.7, 4.65, color=GREEN, width=Pt(2.5))
add_rounded_rect(slide, 4.8, 4.65, 3.8, 0.82,
                 fill=RGBColor(0xE8, 0xF8, 0xF0), line_color=GREEN, line_width=Pt(2))
txt(slide, "Structured JSON Response\n(hubs · links · satellites · hash_definitions · rationale · confidence)",
    4.9, 4.68, 3.6, 0.72, size=Pt(9.5), color=GREEN, align=PP_ALIGN.CENTER, bold=True)

# Token budget note
txt(slide, "📌  Wide tables (>60 cols): top_values and patterns are suppressed to stay within token budget. "
    "max_tokens = 16,000  ·  temperature = 0  (deterministic output)",
    0.3, 5.62, 12.73, 0.42, size=Pt(9), italic=True, color=MID_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — System Prompt Architecture & Modeler Overrides
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
slide_header(slide, "System Prompt Architecture & Modeler Overrides",
             "Standards are stored in the database, versioned, and assembled at runtime — no code change needed to update AI behavior")
footer(slide, 6)

# ── Left half: System Prompt sections ────────────────────────────────────────
txt(slide, "System Prompt — Stored in META.DV_AI_SYSTEM_PROMPT", 0.3, 1.2, 6.5, 0.38,
    size=Pt(13), bold=True, color=NAVY)
txt(slide, "Assembled via LISTAGG(SECTION_CONTENT ORDER BY SECTION_ORDER) at runtime",
    0.3, 1.6, 6.5, 0.3, size=Pt(9.5), italic=True, color=MID_TEXT)

sections = [
    (10, "ROLE",              MID_BLUE,  "Expert DV2.0 banking architect. Goals, available context, output format."),
    (20, "NAMING_CONVENTIONS",TEAL,      "Hub/Link/Sat naming patterns, column roles, UPPER_SNAKE_CASE rules, source system suffix."),
    (30, "HASH_STANDARDS",    SLATE,     "SHA2_BINARY(256), null replacement '-1', UPPER(TRIM), column sort order."),
    (40, "METADATA_COLUMNS",  MID_BLUE,  "Required columns for each entity type: HUB/LNK/SAT/MSAT/ESAT — order and data types."),
    (50, "SATELLITE_RULES",   TEAL,      "One sat per source system. Split by change frequency. FAST/SLOW/STATIC thresholds."),
    (60, "LINK_RULES",        SLATE,     "When to create a link. Alphabetical noun order. Hub reuse detection from registry."),
    (70, "RESPONSE_FORMAT",   ORANGE,    "Strict JSON schema — hubs, links, satellites, hash_definitions, confidence, warnings."),
]

for i, (order, name, col, desc) in enumerate(sections):
    sy = 2.02 + i * 0.68
    # Section order badge
    add_rounded_rect(slide, 0.3, sy, 0.52, 0.52, fill=col, line_color=None)
    txt(slide, str(order), 0.3, sy + 0.08, 0.52, 0.35,
        size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rounded_rect(slide, 0.88, sy, 5.7, 0.52,
                     fill=WHITE, line_color=col, line_width=Pt(1))
    txt(slide, name, 1.0, sy + 0.02, 1.8, 0.26, size=Pt(9.5), bold=True, color=col)
    txt(slide, desc, 2.85, sy + 0.03, 3.65, 0.44, size=Pt(8.5), color=MID_TEXT)

# Divider
add_rect(slide, 6.8, 1.15, 0.04, 6.0, fill=RGBColor(0xCC, 0xD6, 0xE8))

# ── Right half: Modeler Overrides ────────────────────────────────────────────
txt(slide, "Modeler Override Mechanics", 7.0, 1.2, 6.0, 0.38,
    size=Pt(13), bold=True, color=NAVY)

overrides = [
    (AMBER, "Modeler Notes Field",
     "Injected FIRST in the user message, before all other context. "
     "Heading in prompt reads:\n"
     "  === MODELER NOTES — READ THIS FIRST AND APPLY TO EVERY DECISION ===\n"
     "AI is explicitly told: 'Any PK stated here OVERRIDES all statistical PK detection.'"),
    (MID_BLUE, "Business Key Confirmation",
     "Before calling SP_GENERATE_DV_PROPOSAL, modeler reviews PK candidates "
     "scored 0-100. A MODELER_SELECTED row is inserted with PK_SCORE = 100 and "
     "MODELER_SELECTED = TRUE. This propagates to the prompt as '← MODELER CONFIRMED'."),
    (TEAL, "Approved Column Definitions",
     "META.DV_COLUMN_DEFINITIONS stores modeler-approved business definitions "
     "per column. These are injected as a dedicated section:\n"
     "  'APPROVED COLUMN DEFINITIONS — USE THESE VERBATIM'\n"
     "AI cannot paraphrase, shorten, or rewrite them."),
    (GREEN, "Confidence Transparency",
     "AI rates each entity HIGH / MEDIUM / LOW / INFERRED. "
     "LOW or INFERRED confidence triggers a warning in the proposal. "
     "Modeler sees confidence badges in the workbench — prompting them to "
     "add notes and regenerate before approving."),
    (RED, "Deprecation & Exclusion",
     "Modeler notes like 'EXCLUDE COLUMN X — deprecated in v2' cause AI to "
     "omit those columns from all entities. Mention in rationale is mandatory."),
]

for i, (col, title, body) in enumerate(overrides):
    oy = 1.65 + i * 1.1
    add_rounded_rect(slide, 6.95, oy, 6.15, 1.0,
                     fill=WHITE, line_color=col, line_width=Pt(1.5))
    # Left color stripe
    shape = slide.shapes.add_shape(1, Inches(6.95), Inches(oy), Inches(0.1), Inches(1.0))
    shape.fill.solid(); shape.fill.fore_color.rgb = col
    shape.line.fill.background()
    txt(slide, title, 7.15, oy + 0.06, 5.75, 0.28, size=Pt(10.5), bold=True, color=col)
    txt(slide, body, 7.15, oy + 0.35, 5.75, 0.58, size=Pt(8.5), color=MID_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Metadata Tables
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
slide_header(slide, "Metadata Tables — The NEXUS Registry",
             "All platform state lives in META schema — source of truth for AI, approvals, DDL, and dbt generation")
footer(slide, 7)

tables = [
    # (name, color, x, y, description, key_columns)
    ("DV_ENTITY", MID_BLUE, 0.25, 1.25,
     "Central registry of all vault entities",
     "ENTITY_ID · ENTITY_TYPE · APPROVAL_STATUS · DOMAIN"),
    ("DV_ENTITY_COLUMN", MID_BLUE, 0.25, 2.6,
     "Column definitions for each entity",
     "COLUMN_ROLE (HK/BK/ATTR/META) · DATA_TYPE · SOURCE_COLUMN"),
    ("DV_HASH_DEFINITION", TEAL, 4.45, 1.25,
     "Reproducible hash key specifications",
     "ALGORITHM · SOURCE_COLUMNS[] · NULL_REPLACEMENT · DELIMITER"),
    ("DV_SOURCE_MAPPING", TEAL, 4.45, 2.6,
     "Source-to-vault column lineage",
     "SOURCE_TABLE · SOURCE_COLUMN · TARGET_ENTITY · TRANSFORMATION"),
    ("DV_PROFILING_RUN", SLATE, 0.25, 3.95,
     "One row per profiling execution",
     "RUN_ID · ROW_COUNT · PROFILING_METHOD (EXACT/HLL) · STATUS"),
    ("DV_PROFILING_RESULTS", SLATE, 0.25, 5.1,
     "Per-column statistics per run",
     "UNIQUENESS_RATIO · NULL_PCT · CHANGE_FREQUENCY · IS_PK_CANDIDATE"),
    ("DV_PK_CANDIDATES", ORANGE, 4.45, 3.95,
     "Ranked business key candidates",
     "COLUMN_NAMES[] · PK_SCORE (0-100) · MODELER_SELECTED"),
    ("DV_DESIGN_PROPOSAL", ORANGE, 4.45, 5.1,
     "Raw AI-generated proposal JSON",
     "PROPOSAL_JSON (VARIANT) · CONFIDENCE · INPUT_SCENARIO · STATUS"),
    ("DV_DESIGN_WORKSPACE", GREEN, 8.65, 1.25,
     "Modeler editing state + versioning",
     "WORKSPACE_JSON (VARIANT) · VERSION_NUMBER · STATUS · PARENT_WORKSPACE"),
    ("DV_COLUMN_DEFINITIONS", GREEN, 8.65, 2.6,
     "Approved column business definitions",
     "DEFINITION · IS_SENSITIVE · TABLE_DESCRIPTION"),
    ("DV_AI_SYSTEM_PROMPT", AMBER, 8.65, 3.95,
     "Versioned DV2.0 standards (system prompt)",
     "SECTION_NAME · SECTION_ORDER · SECTION_CONTENT · VERSION"),
    ("DV_AUDIT_LOG", RED, 8.65, 5.1,
     "Immutable governance audit trail",
     "ACTION_TYPE · ENTITY_ID · PERFORMED_BY · PERFORMED_AT"),
]

for (name, col, bx, by, desc, cols) in tables:
    add_rounded_rect(slide, bx, by, 4.1, 1.12,
                     fill=WHITE, line_color=col, line_width=Pt(1.5))
    # Top color bar
    shape = slide.shapes.add_shape(5, Inches(bx), Inches(by), Inches(4.1), Inches(0.32))
    shape.adjustments[0] = 0.03
    shape.fill.solid(); shape.fill.fore_color.rgb = col
    shape.line.fill.background()
    txt(slide, name, bx + 0.1, by + 0.02, 3.9, 0.28, size=Pt(9.5), bold=True, color=WHITE)
    txt(slide, desc, bx + 0.1, by + 0.37, 3.9, 0.26, size=Pt(8.5), color=NAVY, bold=False)
    txt(slide, cols, bx + 0.1, by + 0.63, 3.9, 0.4, size=Pt(7.5), italic=True, color=MID_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — dbt Generation (Phase 4)
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
slide_header(slide, "dbt Model Generation — Phase 4",
             "Approved registry metadata drives automated generation of staging, raw vault, and hash models")
footer(slide, 8)

# Flow: Registry → Generator → dbt Models
# Registry box
add_rounded_rect(slide, 0.3, 1.3, 3.5, 4.9,
                 fill=WHITE, line_color=MID_BLUE, line_width=Pt(1.5))
shape = slide.shapes.add_shape(5, Inches(0.3), Inches(1.3), Inches(3.5), Inches(0.52))
shape.adjustments[0] = 0.04
shape.fill.solid(); shape.fill.fore_color.rgb = MID_BLUE
shape.line.fill.background()
txt(slide, "NEXUS Registry (META)", 0.45, 1.33, 3.2, 0.42, size=Pt(11), bold=True, color=WHITE)

reg_feeds = [
    ("DV_ENTITY",          "Vault entity list + type"),
    ("DV_ENTITY_COLUMN",   "Column roles + data types"),
    ("DV_HASH_DEFINITION", "Hash algorithm + source cols"),
    ("DV_SOURCE_MAPPING",  "Source→vault lineage"),
    ("DV_PROFILING_RESULTS","Change freq classification"),
    ("DV_ABBREVIATION",    "Logical name expansion"),
]
for i, (tbl, desc) in enumerate(reg_feeds):
    ry = 1.98 + i * 0.65
    add_rounded_rect(slide, 0.45, ry, 3.2, 0.55,
                     fill=LIGHT_BLUE, line_color=MID_BLUE, line_width=Pt(0.75))
    txt(slide, tbl, 0.58, ry + 0.03, 3.0, 0.25, size=Pt(9), bold=True, color=NAVY)
    txt(slide, desc, 0.58, ry + 0.28, 3.0, 0.22, size=Pt(8), color=MID_TEXT)

# Arrow
add_arrow(slide, 3.8, 3.75, 4.5, 3.75, color=GREEN, width=Pt(3))
txt(slide, "APPROVED\nentities only", 3.62, 3.82, 1.2, 0.4,
    size=Pt(8.5), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Generator box
add_rounded_rect(slide, 4.55, 2.6, 2.7, 2.3,
                 fill=GREEN, line_color=None)
txt(slide, "dbt\nModel\nGenerator\n\n(Streamlit\nPhase 4)", 4.55, 2.65, 2.7, 2.15,
    size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow to models
add_arrow(slide, 7.25, 3.75, 7.75, 3.75, color=GREEN, width=Pt(3))

# dbt model types
dbt_models = [
    ("stg_<source>__<table>.sql", TEAL,
     "Staging model\n"
     "• Cast columns to correct types\n"
     "• Apply UPPER(TRIM) to BK columns\n"
     "• Compute all hash keys (SHA2_BINARY)\n"
     "• Add REC_SRC constant, BATCH_ID\n"
     "• Drive from DV_SOURCE_MAPPING"),
    ("hub_<noun>.sql", MID_BLUE,
     "Hub model\n"
     "• INSERT-IF-NOT-EXISTS pattern\n"
     "• Driven by DV_ENTITY + DV_HASH_DEFINITION\n"
     "• One hub model per approved HUB entity\n"
     "• Primary key = <noun>_HK (BINARY 32)"),
    ("sat_<src>_<noun>_<desc>.sql", ORANGE,
     "Satellite model\n"
     "• HASHDIFF change detection\n"
     "• Insert new record when HASHDIFF changes\n"
     "• Split models for FAST vs SLOW satellites\n"
     "• Column list from DV_ENTITY_COLUMN (ATTR)"),
    ("lnk_<n1>_<n2>.sql", SLATE,
     "Link model\n"
     "• Composite HK from participating hub HKs\n"
     "• Alphabetical hub ordering enforced\n"
     "• Degenerate keys handled as ATTR columns\n"
     "• Source: DV_RELATIONSHIP for FK structure"),
]

for i, (model, col, desc) in enumerate(dbt_models):
    mx = 7.8 + (i % 2) * 2.65
    my = 1.3 + (i // 2) * 2.75
    add_rounded_rect(slide, mx, my, 2.6, 2.55,
                     fill=WHITE, line_color=col, line_width=Pt(1.5))
    shape = slide.shapes.add_shape(5, Inches(mx), Inches(my), Inches(2.6), Inches(0.48))
    shape.adjustments[0] = 0.04
    shape.fill.solid(); shape.fill.fore_color.rgb = col
    shape.line.fill.background()
    txt(slide, model, mx + 0.1, my + 0.05, 2.4, 0.38, size=Pt(8.5), bold=True, color=WHITE)
    txt(slide, desc, mx + 0.1, my + 0.56, 2.4, 1.88, size=Pt(8), color=MID_TEXT)

# Footer note
add_rounded_rect(slide, 0.3, 6.48, 12.73, 0.55,
                 fill=RGBColor(0xE8, 0xF8, 0xF0), line_color=GREEN, line_width=Pt(1))
txt(slide, "📐  dbt models are generated as templated SQL using the APPROVED registry state. "
    "Hash definitions in DV_HASH_DEFINITION ensure bit-for-bit reproducibility across all environments. "
    "Change frequency classification (FAST/SLOW/STATIC from profiling) drives satellite splitting strategy.",
    0.5, 6.52, 12.3, 0.46, size=Pt(9.5), color=RGBColor(0x0A, 0x4D, 0x2F))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Snowflake Native: Pros & Cons
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=OFF_WHITE)
slide_header(slide, "Snowflake Native — Pros & Cons",
             "Building entirely inside Snowflake eliminates external tooling but introduces platform dependencies")
footer(slide, 9)

# Divider
add_rect(slide, 6.55, 1.15, 0.04, 5.95, fill=RGBColor(0xCC, 0xD6, 0xE8))

# ── PROS ──────────────────────────────────────────────────────────────────────
add_rounded_rect(slide, 0.25, 1.18, 6.15, 0.45,
                 fill=GREEN, line_color=None)
txt(slide, "✅  Advantages — Being Snowflake Native", 0.4, 1.22, 5.9, 0.36,
    size=Pt(12), bold=True, color=WHITE)

pros = [
    ("Zero External Infrastructure",
     "No separate AI service, orchestration engine, or UI deployment. "
     "Streamlit in Snowflake, Cortex AI_COMPLETE, and Snowpark SPs are all "
     "natively provisioned. No API keys, no VPCs, no dependency management."),
    ("Data Never Leaves the Platform",
     "Source data, AI prompts, and responses all stay inside Snowflake's "
     "trust boundary. Critical for banking — PII and confidential schema details "
     "are not transmitted to external services."),
    ("Native RBAC & Governance",
     "Snowflake roles control exactly who can call which stored procedure, "
     "see which registry tables, or approve proposals. No separate IAM layer."),
    ("VARIANT for AI Responses",
     "Snowflake's VARIANT type stores AI-generated JSON natively. "
     "Proposals can be queried with dot-notation SQL — no custom serialization."),
    ("Unified Billing",
     "Compute credits cover profiling, AI calls, and UI rendering. "
     "No separate AI vendor invoice or separate BI tool subscription."),
    ("Micro-Partition Pruning",
     "Query performance on registry tables scales automatically. "
     "No custom indexes required — Snowflake's micro-partition pruning handles "
     "predicate pushdown on SOURCE_SYSTEM, STATUS, and date columns."),
    ("Cortex Model Access",
     "New Cortex model versions (Claude Opus 4, 4.5, 4.6…) become available "
     "automatically — just change the model parameter string in the SP, "
     "no infrastructure change needed."),
]

for i, (title, body) in enumerate(pros):
    py = 1.72 + i * 0.72
    add_rounded_rect(slide, 0.25, py, 6.15, 0.67,
                     fill=WHITE, line_color=GREEN, line_width=Pt(0.75))
    # Green left stripe
    shape = slide.shapes.add_shape(1, Inches(0.25), Inches(py), Inches(0.12), Inches(0.67))
    shape.fill.solid(); shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()
    txt(slide, title, 0.47, py + 0.03, 5.75, 0.25, size=Pt(9.5), bold=True, color=GREEN)
    txt(slide, body, 0.47, py + 0.28, 5.75, 0.34, size=Pt(8), color=MID_TEXT)

# ── CONS ──────────────────────────────────────────────────────────────────────
add_rounded_rect(slide, 6.65, 1.18, 6.43, 0.45,
                 fill=RED, line_color=None)
txt(slide, "⚠️  Constraints & Considerations", 6.8, 1.22, 6.15, 0.36,
    size=Pt(12), bold=True, color=WHITE)

cons = [
    ("Snowflake Lock-In",
     "Platform portability is limited. Migrating NEXUS to another cloud data platform "
     "would require rewriting all stored procedures and replacing Cortex with an "
     "equivalent managed AI service."),
    ("Python SP Package Restrictions",
     "Snowpark Python SPs can only use the Snowflake-approved Anaconda channel. "
     "Popular libraries like httpx, pydantic, or sqlmodel are not available without "
     "native app packaging workarounds."),
    ("No Local Development Loop",
     "Every SP change requires a Snowflake connection to test. There is no local "
     "unit-test framework for Snowpark procedures — debugging requires repeated "
     "cloud deploys."),
    ("Cortex Model Constraints",
     "AI model selection is limited to models supported by SNOWFLAKE.CORTEX.AI_COMPLETE. "
     "Fine-tuning is not supported. Context window and token rate limits are controlled "
     "by Snowflake, not the team."),
    ("Streamlit UI Limitations",
     "Streamlit in Snowflake has a smaller component ecosystem than open-source Streamlit. "
     "Custom React components, websockets, and some third-party libraries are unavailable "
     "or require workarounds."),
    ("Compute Cost Visibility",
     "AI_COMPLETE calls consume Snowflake credits alongside query compute — it can be "
     "difficult to isolate AI spend from pipeline spend without dedicated warehouses "
     "per workload."),
    ("Latency on First AI Call",
     "Cortex cold-start on first call in a session can add latency. Proposal generation "
     "for wide tables (>60 cols, 16K token output) can take 15–40 seconds — "
     "async patterns not natively available in Streamlit SiS."),
]

for i, (title, body) in enumerate(cons):
    cy = 1.72 + i * 0.72
    add_rounded_rect(slide, 6.65, cy, 6.43, 0.67,
                     fill=WHITE, line_color=RED, line_width=Pt(0.75))
    shape = slide.shapes.add_shape(1, Inches(6.65), Inches(cy), Inches(0.12), Inches(0.67))
    shape.fill.solid(); shape.fill.fore_color.rgb = RED
    shape.line.fill.background()
    txt(slide, title, 6.87, cy + 0.03, 6.05, 0.25, size=Pt(9.5), bold=True, color=RED)
    txt(slide, body, 6.87, cy + 0.28, 6.05, 0.34, size=Pt(8), color=MID_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Summary & Roadmap
# ══════════════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(slide, 0, 0, 0.18, 7.5, fill=TEAL)
add_rect(slide, 0.18, 2.1, 13.15, 0.06, fill=TEAL)

txt(slide, "Where We Are & What's Next", 0.5, 0.4, 12, 0.65,
    size=Pt(28), bold=True, color=WHITE)

phases = [
    ("Phase 1\nFOUNDATION", GREEN, "COMPLETE",
     "Database & schemas · 12 registry tables\n"
     "Abbreviation table (Erwin .ABR sync)\n"
     "Versioned system prompt · Reference hubs"),
    ("Phase 2\nPROFILING", GREEN, "COMPLETE",
     "SP_PROFILE_TABLE (EXACT + HLL)\n"
     "PK candidate scoring & detection\n"
     "Change frequency classification"),
    ("Phase 3\nAI WORKBENCH", AMBER, "IN PROGRESS",
     "SP_GENERATE_DV_PROPOSAL · Claude Opus 4.6\n"
     "Design Workbench (Streamlit)\n"
     "Approve → Registry → Audit log"),
    ("Phase 4\ndbt GENERATION", SLATE, "PLANNED",
     "Generate stg / hub / sat / lnk models\n"
     "Hash key templates from registry\n"
     "Satellite split by change frequency"),
]

for i, (phase, col, status, desc) in enumerate(phases):
    px = 0.55 + i * 3.15
    add_rounded_rect(slide, px, 2.35, 3.0, 3.0,
                     fill=RGBColor(0x1A, 0x38, 0x66), line_color=col, line_width=Pt(2))
    # Status badge
    add_rounded_rect(slide, px + 0.15, 2.42, 1.5, 0.38,
                     fill=col, line_color=None)
    txt(slide, status, px + 0.15, 2.44, 1.5, 0.32,
        size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, phase, px + 0.15, 2.87, 2.7, 0.65,
        size=Pt(12), bold=True, color=col)
    txt(slide, desc, px + 0.15, 3.6, 2.75, 1.65,
        size=Pt(9.5), color=RGBColor(0xA8, 0xC7, 0xE8))

# Key stats
txt(slide, "Platform at a Glance", 0.5, 5.55, 6, 0.38,
    size=Pt(13), bold=True, color=TEAL)

stats = [
    ("12", "META Registry Tables"),
    ("7",  "System Prompt Sections"),
    ("4",  "Input Scenarios\n(Full Profile → Inference)"),
    ("5",  "Context Layers\nper AI Call"),
    ("100","PK Score for Modeler\nConfirmed Keys"),
]

for i, (num_val, label) in enumerate(stats):
    sx = 0.55 + i * 2.55
    txt(slide, num_val, sx, 5.95, 1.5, 0.7,
        size=Pt(40), bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    txt(slide, label, sx, 6.65, 2.3, 0.45,
        size=Pt(9), color=RGBColor(0xA8, 0xC7, 0xE8), align=PP_ALIGN.CENTER)

# Bottom bar
add_rect(slide, 0, 7.25, 13.33, 0.25, fill=RGBColor(0x0A, 0x16, 0x2B))
txt(slide, "NEXUS DV2.0  |  AI-Assisted Data Vault Automation  |  Built on Snowflake Cortex",
    0.3, 7.27, 12, 0.2, size=Pt(8), color=RGBColor(0x80, 0x9A, 0xBF))


# ── Save ──────────────────────────────────────────────────────────────────────
output_path = r"C:\Users\Vijay RK\Documents\Projects\AIDataArchitect\NEXUS_DV2_Platform_Overview.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
